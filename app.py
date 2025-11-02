import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime, timedelta
import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import tempfile
import os

# ページ設定
st.set_page_config(
    page_title="不良分析QC7つ道具システム",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ローカル完結設定：PlotlyのCDNを無効化
import streamlit.components.v1 as components
config = {
    'displayModeBar': True,
    'displaylogo': False,
    'modeBarButtonsToRemove': ['pan2d', 'lasso2d'],
    'toImageButtonOptions': {
        'format': 'png',
        'filename': 'qc7_graph',
        'height': 600,
        'width': 1000,
        'scale': 1
    }
}
# plotly.ioのデフォルト設定
try:
    import plotly.io as pio
    pio.renderers.default = "browser"  # ローカルレンダリング
except:
    pass

# 日本語フォント設定
plt.rcParams['font.family'] = 'DejaVu Sans'

# QC7つ道具の実装関数
def create_pareto_chart(df):
    """パレート図の作成"""
    st.subheader("📊 パレート図 - 不良項目別分析")
    
    if '不良項目' not in df.columns or '不良数' not in df.columns:
        st.error("データに「不良項目」と「不良数」の列が必要です。")
        return
    
    # 不良項目別の集計
    defect_summary = df.groupby('不良項目')['不良数'].sum().sort_values(ascending=False)
    
    # 累積比率の計算
    cumulative_ratio = (defect_summary.cumsum() / defect_summary.sum() * 100).round(1)
    
    # パレート図の作成
    fig = make_subplots(specs=[[{"secondary_y": True}]])
    
    # 棒グラフ
    fig.add_trace(
        go.Bar(x=defect_summary.index, y=defect_summary.values, name="不良数", marker_color='lightblue'),
        secondary_y=False,
    )
    
    # 累積比率線
    fig.add_trace(
        go.Scatter(x=defect_summary.index, y=cumulative_ratio.values, 
                  mode='lines+markers', name="累積比率", line=dict(color='red', width=3)),
        secondary_y=True,
    )
    
    # 80%ライン
    fig.add_hline(y=80, line_dash="dash", line_color="red", 
                  annotation_text="80%ライン", secondary_y=True)
    
    # レイアウト設定
    fig.update_xaxes(title_text="不良項目")
    fig.update_yaxes(title_text="不良数", secondary_y=False)
    fig.update_yaxes(title_text="累積比率 (%)", secondary_y=True)
    fig.update_layout(title_text="パレート図 - 不良項目別分析", height=500)
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 画像ダウンロードボタン
    img_bytes = fig.to_image(format="png", width=1200, height=600)
    st.download_button(
        label="📥 PNG画像をダウンロード",
        data=img_bytes,
        file_name=f"パレート図_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png",
        mime="image/png"
    )
    
    # 分析結果
    st.subheader("📈 分析結果")
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**不良項目ランキング**")
        for i, (item, count) in enumerate(defect_summary.items(), 1):
            ratio = (count / defect_summary.sum() * 100)
            st.write(f"{i}. {item}: {count}件 ({ratio:.1f}%)")
    
    with col2:
        st.write("**重点管理項目（上位80%）**")
        top_80_items = defect_summary[cumulative_ratio <= 80]
        for item in top_80_items.index:
            st.write(f"• {item}")

def create_fishbone_diagram(df):
    """特性要因図の作成"""
    st.subheader("🐟 特性要因図 - 不良原因分析")
    
    if '原因分類' not in df.columns:
        st.error("データに「原因分類」の列が必要です。")
        return
    
    # 原因別の集計
    cause_summary = df.groupby('原因分類')['不良数'].sum().sort_values(ascending=False)
    
    # 4M分類
    m4_categories = {
        'Man': ['作業者', '人', 'オペレーター'],
        'Machine': ['機械', '設備', '工具', '加工'],
        'Material': ['材料', '部品', '素材'],
        'Method': ['方法', '手順', '環境', '条件']
    }
    
    # 原因を4Mに分類
    categorized_causes = {}
    for category, keywords in m4_categories.items():
        categorized_causes[category] = []
        for cause in cause_summary.index:
            if any(keyword in cause for keyword in keywords):
                categorized_causes[category].append((cause, cause_summary[cause]))
    
    # 特性要因図の表示
    st.write("**4M分析による原因分類**")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**Man（人）**")
        for cause, count in categorized_causes['Man']:
            st.write(f"• {cause}: {count}件")
        
        st.write("**Machine（機械）**")
        for cause, count in categorized_causes['Machine']:
            st.write(f"• {cause}: {count}件")
    
    with col2:
        st.write("**Material（材料）**")
        for cause, count in categorized_causes['Material']:
            st.write(f"• {cause}: {count}件")
        
        st.write("**Method（方法）**")
        for cause, count in categorized_causes['Method']:
            st.write(f"• {cause}: {count}件")
    
    # 原因別円グラフ
    fig = px.pie(values=cause_summary.values, names=cause_summary.index, 
                 title="原因分類別不良数")
    st.plotly_chart(fig, use_container_width=True)

def create_histogram(df):
    """ヒストグラムの作成"""
    st.subheader("📊 ヒストグラム - 不良率分布分析")
    
    if '不良数' not in df.columns or '検査数' not in df.columns:
        st.error("データに「不良数」と「検査数」の列が必要です。")
        return
    
    # 不良率の計算
    df['不良率'] = (df['不良数'] / df['検査数'] * 100).round(2)
    
    # ヒストグラムの作成
    fig = px.histogram(df, x='不良率', nbins=20, 
                       title="不良率分布ヒストグラム",
                       labels={'不良率': '不良率 (%)', 'count': '頻度'})
    
    # 平均線の追加
    mean_rate = df['不良率'].mean()
    fig.add_vline(x=mean_rate, line_dash="dash", line_color="red",
                  annotation_text=f"平均: {mean_rate:.2f}%")
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 統計情報
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.metric("平均不良率", f"{df['不良率'].mean():.2f}%")
    
    with col2:
        st.metric("標準偏差", f"{df['不良率'].std():.2f}%")
    
    with col3:
        st.metric("最大不良率", f"{df['不良率'].max():.2f}%")

def create_scatter_plot(df):
    """散布図の作成"""
    st.subheader("📈 散布図 - 相関分析")
    
    if '不良数' not in df.columns or '検査数' not in df.columns:
        st.error("データに「不良数」と「検査数」の列が必要です。")
        return
    
    # 不良率の計算
    df['不良率'] = (df['不良数'] / df['検査数'] * 100).round(2)
    
    # 散布図の選択肢
    st.write("**分析項目を選択してください**")
    
    col1, col2 = st.columns(2)
    
    with col1:
        x_axis = st.selectbox("X軸", ['検査数', '不良数', '不良率'])
    
    with col2:
        y_axis = st.selectbox("Y軸", ['不良率', '不良数', '検査数'])
    
    # 散布図の作成
    fig = px.scatter(df, x=x_axis, y=y_axis, 
                     color='原因分類' if '原因分類' in df.columns else None,
                     title=f"{x_axis} vs {y_axis} 散布図",
                     hover_data=['製品名', '不良項目'] if '製品名' in df.columns else None)
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 相関係数の計算
    correlation = df[x_axis].corr(df[y_axis])
    st.write(f"**相関係数**: {correlation:.3f}")
    
    if abs(correlation) > 0.7:
        st.success("強い相関関係があります")
    elif abs(correlation) > 0.3:
        st.warning("中程度の相関関係があります")
    else:
        st.info("弱い相関関係です")

def create_control_chart(df):
    """管理図の作成"""
    st.subheader("📊 管理図 - 工程管理")
    
    if '日付' not in df.columns or '不良数' not in df.columns or '検査数' not in df.columns:
        st.error("データに「日付」「不良数」「検査数」の列が必要です。")
        return
    
    # 日付でソート
    df_sorted = df.sort_values('日付')
    
    # 日付別の集計
    daily_summary = df_sorted.groupby('日付').agg({
        '不良数': 'sum',
        '検査数': 'sum'
    }).reset_index()
    
    # 不良率の計算
    daily_summary['不良率'] = (daily_summary['不良数'] / daily_summary['検査数'] * 100).round(2)
    
    # 管理限界の計算（p管理図）
    p_bar = daily_summary['不良率'].mean()
    n_bar = daily_summary['検査数'].mean()
    
    # 3σ管理限界
    ucl = p_bar + 3 * np.sqrt(p_bar * (100 - p_bar) / n_bar)
    lcl = max(0, p_bar - 3 * np.sqrt(p_bar * (100 - p_bar) / n_bar))
    
    # 管理図の作成
    fig = go.Figure()
    
    # データ点
    fig.add_trace(go.Scatter(
        x=daily_summary['日付'],
        y=daily_summary['不良率'],
        mode='lines+markers',
        name='不良率',
        line=dict(color='blue')
    ))
    
    # 中心線
    fig.add_hline(y=p_bar, line_dash="dash", line_color="green",
                  annotation_text=f"中心線: {p_bar:.2f}%")
    
    # 管理限界
    fig.add_hline(y=ucl, line_dash="dash", line_color="red",
                  annotation_text=f"UCL: {ucl:.2f}%")
    fig.add_hline(y=lcl, line_dash="dash", line_color="red",
                  annotation_text=f"LCL: {lcl:.2f}%")
    
    fig.update_layout(
        title="p管理図（不良率管理図）",
        xaxis_title="日付",
        yaxis_title="不良率 (%)",
        height=500
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 異常値の検出
    outliers = daily_summary[(daily_summary['不良率'] > ucl) | (daily_summary['不良率'] < lcl)]
    
    if len(outliers) > 0:
        st.warning(f"⚠️ {len(outliers)}件の異常値が検出されました")
        st.dataframe(outliers)
    else:
        st.success("✅ 異常値は検出されませんでした")

def create_checklist(df):
    """チェックシートの作成"""
    st.subheader("📋 チェックシート - データ収集支援")
    
    # 不良項目別のチェックリスト
    if '不良項目' in df.columns:
        st.write("**不良項目別チェックリスト**")
        
        defect_items = df['不良項目'].unique()
        
        for item in defect_items:
            with st.expander(f"📌 {item}"):
                item_data = df[df['不良項目'] == item]
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.write(f"**発生件数**: {len(item_data)}件")
                    st.write(f"**総不良数**: {item_data['不良数'].sum()}件")
                
                with col2:
                    if '原因分類' in df.columns:
                        causes = item_data['原因分類'].value_counts()
                        st.write("**主な原因**:")
                        for cause, count in causes.head(3).items():
                            st.write(f"• {cause}: {count}件")
    
    # カスタムチェックリスト
    st.write("**カスタムチェックリスト**")
    
    checklist_items = st.text_area(
        "チェック項目を入力してください（1行に1項目）",
        value="寸法測定\n表面粗さ検査\n外観検査\n機能テスト\n包装確認",
        height=100
    )
    
    if st.button("チェックリストを生成"):
        items = [item.strip() for item in checklist_items.split('\n') if item.strip()]
        
        st.write("**生成されたチェックシート**")
        for i, item in enumerate(items, 1):
            st.checkbox(f"{i}. {item}", key=f"check_{i}")

def create_graphs(df):
    """グラフの作成"""
    st.subheader("📊 グラフ - 時系列・比較分析")
    
    # グラフタイプの選択
    graph_type = st.selectbox(
        "グラフタイプを選択してください",
        ["時系列グラフ", "工程別比較", "製品別比較", "原因別比較"]
    )
    
    if graph_type == "時系列グラフ":
        if '日付' not in df.columns:
            st.error("データに「日付」の列が必要です。")
            return
        
        # 日付別の集計
        daily_summary = df.groupby('日付').agg({
            '不良数': 'sum',
            '検査数': 'sum'
        }).reset_index()
        daily_summary['不良率'] = (daily_summary['不良数'] / daily_summary['検査数'] * 100).round(2)
        
        # 時系列グラフ
        fig = px.line(daily_summary, x='日付', y='不良率',
                      title="不良率推移（時系列）")
        st.plotly_chart(fig, use_container_width=True)
    
    elif graph_type == "工程別比較":
        if '発生工程' not in df.columns:
            st.error("データに「発生工程」の列が必要です。")
            return
        
        # 工程別の集計
        process_summary = df.groupby('発生工程').agg({
            '不良数': 'sum',
            '検査数': 'sum'
        }).reset_index()
        process_summary['不良率'] = (process_summary['不良数'] / process_summary['検査数'] * 100).round(2)
        
        # 工程別比較グラフ
        fig = px.bar(process_summary, x='発生工程', y='不良率',
                     title="工程別不良率比較")
        st.plotly_chart(fig, use_container_width=True)
    
    elif graph_type == "製品別比較":
        if '製品名' not in df.columns:
            st.error("データに「製品名」の列が必要です。")
            return
        
        # 製品別の集計
        product_summary = df.groupby('製品名').agg({
            '不良数': 'sum',
            '検査数': 'sum'
        }).reset_index()
        product_summary['不良率'] = (product_summary['不良数'] / product_summary['検査数'] * 100).round(2)
        
        # 製品別比較グラフ
        fig = px.bar(product_summary, x='製品名', y='不良率',
                     title="製品別不良率比較")
        st.plotly_chart(fig, use_container_width=True)
    
    elif graph_type == "原因別比較":
        if '原因分類' not in df.columns:
            st.error("データに「原因分類」の列が必要です。")
            return
        
        # 原因別の集計
        cause_summary = df.groupby('原因分類')['不良数'].sum().sort_values(ascending=False)
        
        # 原因別比較グラフ
        fig = px.pie(values=cause_summary.values, names=cause_summary.index,
                     title="原因分類別不良数")
        st.plotly_chart(fig, use_container_width=True)

# カスタムCSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
    }
    .tool-card {
        border: 1px solid #ddd;
        border-radius: 10px;
        padding: 1rem;
        margin: 0.5rem 0;
        background-color: #f9f9f9;
    }
    .metric-card {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 5px;
        margin: 0.5rem 0;
    }
</style>
""", unsafe_allow_html=True)

# メインタイトル
st.markdown('<h1 class="main-header">📊 不良分析QC7つ道具システム</h1>', unsafe_allow_html=True)

# サイドバー
st.sidebar.title("🔧 メニュー")

# データアップロード
st.sidebar.header("📁 データアップロード")
uploaded_file = st.sidebar.file_uploader(
    "CSVまたはExcelファイルを選択してください",
    type=['csv', 'xlsx', 'xls'],
    help="不良データのCSVまたはExcelファイルをアップロードしてください"
)

# データ入力機能
st.sidebar.subheader("📝 データ入力")
if st.sidebar.button("📝 手動データ入力"):
    st.session_state.show_data_input = True

# 手動データ入力画面
if st.session_state.get('show_data_input', False):
    st.header("📝 データ入力")
    
    # 入力フォーム
    with st.form("data_input_form"):
        col1, col2 = st.columns(2)
        
        with col1:
            input_date = st.date_input("日付", value=datetime.now().date())
            product_name = st.selectbox("製品名", ["製品A", "製品B", "製品C", "製品D"])
            defect_type = st.selectbox("不良項目", ["寸法不良", "表面粗さ", "欠け", "傷", "変形"])
            defect_count = st.number_input("不良数", min_value=0, value=0)
        
        with col2:
            inspection_count = st.number_input("検査数", min_value=1, value=100)
            cause_category = st.selectbox("原因分類", ["加工", "工具", "材料", "作業者", "環境"])
            process = st.selectbox("発生工程", ["旋盤", "研削", "組立", "熱処理"])
            remarks = st.text_input("備考", placeholder="詳細な原因など")
        
        submitted = st.form_submit_button("データを追加")
        
        if submitted:
            # 新しいデータを追加
            new_data = {
                '日付': input_date.strftime('%Y-%m-%d'),
                '製品名': product_name,
                '不良項目': defect_type,
                '不良数': defect_count,
                '検査数': inspection_count,
                '原因分類': cause_category,
                '発生工程': process,
                '備考': remarks
            }
            
            # セッション状態にデータを追加
            if 'manual_data' not in st.session_state:
                st.session_state.manual_data = []
            
            st.session_state.manual_data.append(new_data)
            st.success("データを追加しました！")
    
    # 入力済みデータの表示
    if 'manual_data' in st.session_state and st.session_state.manual_data:
        st.subheader("📋 入力済みデータ")
        manual_df = pd.DataFrame(st.session_state.manual_data)
        st.dataframe(manual_df, use_container_width=True)
        
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("✅ データを確定"):
                st.session_state.df = manual_df
                st.session_state.show_data_input = False
                st.success("データが確定されました！")
                st.rerun()
        
        with col2:
            if st.button("🗑️ データをクリア"):
                st.session_state.manual_data = []
                st.success("データをクリアしました！")
                st.rerun()
        
        with col3:
            if st.button("📥 CSVダウンロード"):
                csv = manual_df.to_csv(index=False, encoding='utf-8')
                st.download_button(
                    label="CSVファイルをダウンロード",
                    data=csv,
                    file_name=f"不良データ_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                    mime="text/csv"
                )
    
    # 一括データ入力
    st.subheader("📊 一括データ入力")
    st.write("複数のデータを一度に入力できます")
    
    with st.form("bulk_input_form"):
        bulk_data = st.text_area(
            "データを入力してください（1行に1件、カンマ区切り）",
            placeholder="日付,製品名,不良項目,不良数,検査数,原因分類,発生工程,備考\n2024-01-01,製品A,寸法不良,5,100,加工,旋盤,工具摩耗",
            height=200
        )
        
        if st.form_submit_button("一括データを追加"):
            try:
                lines = bulk_data.strip().split('\n')
                for line in lines:
                    if line.strip():
                        parts = [part.strip() for part in line.split(',')]
                        if len(parts) >= 8:
                            bulk_record = {
                                '日付': parts[0],
                                '製品名': parts[1],
                                '不良項目': parts[2],
                                '不良数': int(parts[3]),
                                '検査数': int(parts[4]),
                                '原因分類': parts[5],
                                '発生工程': parts[6],
                                '備考': parts[7] if len(parts) > 7 else ''
                            }
                            
                            if 'manual_data' not in st.session_state:
                                st.session_state.manual_data = []
                            
                            st.session_state.manual_data.append(bulk_record)
                
                st.success(f"{len(lines)}件のデータを追加しました！")
                st.rerun()
                
            except Exception as e:
                st.error(f"データの形式が正しくありません: {str(e)}")
                st.info("正しい形式: 日付,製品名,不良項目,不良数,検査数,原因分類,発生工程,備考")

# サンプルデータ生成ボタン
if st.sidebar.button("📋 サンプルデータを生成"):
    # サンプルデータ生成
    np.random.seed(42)
    dates = pd.date_range(start='2024-01-01', end='2024-01-31', freq='D')
    
    sample_data = []
    products = ['製品A', '製品B', '製品C']
    defect_types = ['寸法不良', '表面粗さ', '欠け', '傷', '変形']
    causes = ['加工', '材料', '工具', '作業者', '環境']
    processes = ['旋盤', 'フライス', '研削', '組立', '検査']
    
    for date in dates:
        for product in products:
            for _ in range(np.random.randint(1, 4)):
                defect_type = np.random.choice(defect_types)
                defect_count = np.random.randint(1, 10)
                inspection_count = np.random.randint(80, 120)
                cause = np.random.choice(causes)
                process = np.random.choice(processes)
                
                sample_data.append({
                    '日付': date.strftime('%Y-%m-%d'),
                    '製品名': product,
                    '不良項目': defect_type,
                    '不良数': defect_count,
                    '検査数': inspection_count,
                    '原因分類': cause,
                    '発生工程': process,
                    '備考': f'{cause}による{defect_type}'
                })
    
    df = pd.DataFrame(sample_data)
    st.session_state.df = df
    st.sidebar.success("サンプルデータを生成しました！")

# データ読み込み
df = None
if uploaded_file is not None:
    try:
        if uploaded_file.name.endswith('.csv'):
            df = pd.read_csv(uploaded_file, encoding='utf-8')
        else:
            df = pd.read_excel(uploaded_file)
        
        # データをセッション状態に保存
        st.session_state.df = df
        st.sidebar.success(f"データを読み込みました！ ({len(df)}件)")
        
    except Exception as e:
        st.sidebar.error(f"ファイル読み込みエラー: {str(e)}")

# セッション状態からデータを取得
if 'df' in st.session_state:
    df = st.session_state.df
    
    # データ概要表示
    st.header("📈 データ概要")
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("総データ数", len(df))
    
    with col2:
        total_defects = df['不良数'].sum() if '不良数' in df.columns else 0
        st.metric("総不良数", total_defects)
    
    with col3:
        total_inspections = df['検査数'].sum() if '検査数' in df.columns else 0
        st.metric("総検査数", total_inspections)
    
    with col4:
        defect_rate = (total_defects / total_inspections * 100) if total_inspections > 0 else 0
        st.metric("不良率", f"{defect_rate:.2f}%")
    
    # データプレビュー
    st.subheader("📋 データプレビュー")
    st.dataframe(df.head(10), use_container_width=True)
    
    # QC7つ道具選択
    st.sidebar.header("🔧 QC7つ道具")
    
    tools = {
        "1. パレート図": "pareto",
        "2. 特性要因図": "fishbone", 
        "3. ヒストグラム": "histogram",
        "4. 散布図": "scatter",
        "5. 管理図": "control",
        "6. チェックシート": "checklist",
        "7. グラフ": "graphs",
        "8. プレゼン資料生成": "presentation"
    }
    
    selected_tool = st.sidebar.selectbox(
        "分析ツールを選択してください",
        list(tools.keys())
    )
    
    # メインコンテンツエリア
    st.header(f"🔍 {selected_tool}")
    
    # 各ツールの実装
    if tools[selected_tool] == "pareto":
        create_pareto_chart(df)
    elif tools[selected_tool] == "fishbone":
        create_fishbone_diagram(df)
    elif tools[selected_tool] == "histogram":
        create_histogram(df)
    elif tools[selected_tool] == "scatter":
        create_scatter_plot(df)
    elif tools[selected_tool] == "control":
        create_control_chart(df)
    elif tools[selected_tool] == "checklist":
        create_checklist(df)
    elif tools[selected_tool] == "graphs":
        create_graphs(df)
    elif tools[selected_tool] == "presentation":
        generate_presentation(df)

else:
    # データ未アップロード時の表示
    st.info("👆 サイドバーからデータファイルをアップロードするか、サンプルデータを生成してください。")
    
    # システム説明
    st.markdown("""
    ## 🎯 システム概要
    
    このシステムは、中小工場の不良分析に特化したQC7つ道具を自動生成します。
    
    ### 📊 対応するQC7つ道具
    1. **パレート図** - 不良項目の重要度を視覚化
    2. **特性要因図** - 不良原因の体系的整理
    3. **ヒストグラム** - 不良率の分布分析
    4. **散布図** - 工程パラメータとの相関分析
    5. **管理図** - 工程の安定性監視
    6. **チェックシート** - データ収集の標準化
    7. **グラフ** - 時系列・比較分析
    
    ### 📁 対応データ形式
    - CSVファイル（UTF-8エンコーディング）
    - Excelファイル（.xlsx, .xls）
    
    ### 📋 必要なデータ項目
    - 日付
    - 製品名
    - 不良項目
    - 不良数
    - 検査数
    - 原因分類
    - 発生工程
    - 備考（任意）
    """)

# QC7つ道具の実装関数
def create_pareto_chart(df):
    """パレート図の作成"""
    st.subheader("📊 パレート図 - 不良項目別分析")
    
    if '不良項目' not in df.columns or '不良数' not in df.columns:
        st.error("データに「不良項目」と「不良数」の列が必要です。")
        return
    
    # 不良項目別の集計
    defect_summary = df.groupby('不良項目')['不良数'].sum().sort_values(ascending=False)
    
    # 累積比率の計算
    cumulative_ratio = (defect_summary.cumsum() / defect_summary.sum() * 100).round(1)
    
    # パレート図の作成
    fig = make_subplots(specs=[[{"secondary_y": True}]])
    
    # 棒グラフ
    fig.add_trace(
        go.Bar(x=defect_summary.index, y=defect_summary.values, name="不良数", marker_color='lightblue'),
        secondary_y=False,
    )
    
    # 累積比率線
    fig.add_trace(
        go.Scatter(x=defect_summary.index, y=cumulative_ratio.values, 
                  mode='lines+markers', name="累積比率", line=dict(color='red', width=3)),
        secondary_y=True,
    )
    
    # 80%ライン
    fig.add_hline(y=80, line_dash="dash", line_color="red", 
                  annotation_text="80%ライン", secondary_y=True)
    
    # レイアウト設定
    fig.update_xaxes(title_text="不良項目")
    fig.update_yaxes(title_text="不良数", secondary_y=False)
    fig.update_yaxes(title_text="累積比率 (%)", secondary_y=True)
    fig.update_layout(title_text="パレート図 - 不良項目別分析", height=500)
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 画像ダウンロードボタン
    img_bytes = fig.to_image(format="png", width=1200, height=600)
    st.download_button(
        label="📥 PNG画像をダウンロード",
        data=img_bytes,
        file_name=f"パレート図_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png",
        mime="image/png"
    )
    
    # 分析結果
    st.subheader("📈 分析結果")
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**不良項目ランキング**")
        for i, (item, count) in enumerate(defect_summary.items(), 1):
            ratio = (count / defect_summary.sum() * 100)
            st.write(f"{i}. {item}: {count}件 ({ratio:.1f}%)")
    
    with col2:
        st.write("**重点管理項目（上位80%）**")
        top_80_items = defect_summary[cumulative_ratio <= 80]
        for item in top_80_items.index:
            st.write(f"• {item}")

def create_fishbone_diagram(df):
    """特性要因図の作成"""
    st.subheader("🐟 特性要因図 - 不良原因分析")
    
    if '原因分類' not in df.columns:
        st.error("データに「原因分類」の列が必要です。")
        return
    
    # 原因別の集計
    cause_summary = df.groupby('原因分類')['不良数'].sum().sort_values(ascending=False)
    
    # 4M分類
    m4_categories = {
        'Man': ['作業者', '人', 'オペレーター'],
        'Machine': ['機械', '設備', '工具', '加工'],
        'Material': ['材料', '部品', '素材'],
        'Method': ['方法', '手順', '環境', '条件']
    }
    
    # 原因を4Mに分類
    categorized_causes = {}
    for category, keywords in m4_categories.items():
        categorized_causes[category] = []
        for cause in cause_summary.index:
            if any(keyword in cause for keyword in keywords):
                categorized_causes[category].append((cause, cause_summary[cause]))
    
    # 特性要因図の表示
    st.write("**4M分析による原因分類**")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**Man（人）**")
        for cause, count in categorized_causes['Man']:
            st.write(f"• {cause}: {count}件")
        
        st.write("**Machine（機械）**")
        for cause, count in categorized_causes['Machine']:
            st.write(f"• {cause}: {count}件")
    
    with col2:
        st.write("**Material（材料）**")
        for cause, count in categorized_causes['Material']:
            st.write(f"• {cause}: {count}件")
        
        st.write("**Method（方法）**")
        for cause, count in categorized_causes['Method']:
            st.write(f"• {cause}: {count}件")
    
    # 原因別円グラフ
    fig = px.pie(values=cause_summary.values, names=cause_summary.index, 
                 title="原因分類別不良数")
    st.plotly_chart(fig, use_container_width=True)

def create_histogram(df):
    """ヒストグラムの作成"""
    st.subheader("📊 ヒストグラム - 不良率分布分析")
    
    if '不良数' not in df.columns or '検査数' not in df.columns:
        st.error("データに「不良数」と「検査数」の列が必要です。")
        return
    
    # 不良率の計算
    df['不良率'] = (df['不良数'] / df['検査数'] * 100).round(2)
    
    # ヒストグラムの作成
    fig = px.histogram(df, x='不良率', nbins=20, 
                       title="不良率分布ヒストグラム",
                       labels={'不良率': '不良率 (%)', 'count': '頻度'})
    
    # 平均線の追加
    mean_rate = df['不良率'].mean()
    fig.add_vline(x=mean_rate, line_dash="dash", line_color="red",
                  annotation_text=f"平均: {mean_rate:.2f}%")
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 統計情報
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.metric("平均不良率", f"{df['不良率'].mean():.2f}%")
    
    with col2:
        st.metric("標準偏差", f"{df['不良率'].std():.2f}%")
    
    with col3:
        st.metric("最大不良率", f"{df['不良率'].max():.2f}%")

def create_scatter_plot(df):
    """散布図の作成"""
    st.subheader("📈 散布図 - 相関分析")
    
    if '不良数' not in df.columns or '検査数' not in df.columns:
        st.error("データに「不良数」と「検査数」の列が必要です。")
        return
    
    # 不良率の計算
    df['不良率'] = (df['不良数'] / df['検査数'] * 100).round(2)
    
    # 散布図の選択肢
    st.write("**分析項目を選択してください**")
    
    col1, col2 = st.columns(2)
    
    with col1:
        x_axis = st.selectbox("X軸", ['検査数', '不良数', '不良率'])
    
    with col2:
        y_axis = st.selectbox("Y軸", ['不良率', '不良数', '検査数'])
    
    # 散布図の作成
    fig = px.scatter(df, x=x_axis, y=y_axis, 
                     color='原因分類' if '原因分類' in df.columns else None,
                     title=f"{x_axis} vs {y_axis} 散布図",
                     hover_data=['製品名', '不良項目'] if '製品名' in df.columns else None)
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 相関係数の計算
    correlation = df[x_axis].corr(df[y_axis])
    st.write(f"**相関係数**: {correlation:.3f}")
    
    if abs(correlation) > 0.7:
        st.success("強い相関関係があります")
    elif abs(correlation) > 0.3:
        st.warning("中程度の相関関係があります")
    else:
        st.info("弱い相関関係です")

def create_control_chart(df):
    """管理図の作成"""
    st.subheader("📊 管理図 - 工程管理")
    
    if '日付' not in df.columns or '不良数' not in df.columns or '検査数' not in df.columns:
        st.error("データに「日付」「不良数」「検査数」の列が必要です。")
        return
    
    # 日付でソート
    df_sorted = df.sort_values('日付')
    
    # 日付別の集計
    daily_summary = df_sorted.groupby('日付').agg({
        '不良数': 'sum',
        '検査数': 'sum'
    }).reset_index()
    
    # 不良率の計算
    daily_summary['不良率'] = (daily_summary['不良数'] / daily_summary['検査数'] * 100).round(2)
    
    # 管理限界の計算（p管理図）
    p_bar = daily_summary['不良率'].mean()
    n_bar = daily_summary['検査数'].mean()
    
    # 3σ管理限界
    ucl = p_bar + 3 * np.sqrt(p_bar * (100 - p_bar) / n_bar)
    lcl = max(0, p_bar - 3 * np.sqrt(p_bar * (100 - p_bar) / n_bar))
    
    # 管理図の作成
    fig = go.Figure()
    
    # データ点
    fig.add_trace(go.Scatter(
        x=daily_summary['日付'],
        y=daily_summary['不良率'],
        mode='lines+markers',
        name='不良率',
        line=dict(color='blue')
    ))
    
    # 中心線
    fig.add_hline(y=p_bar, line_dash="dash", line_color="green",
                  annotation_text=f"中心線: {p_bar:.2f}%")
    
    # 管理限界
    fig.add_hline(y=ucl, line_dash="dash", line_color="red",
                  annotation_text=f"UCL: {ucl:.2f}%")
    fig.add_hline(y=lcl, line_dash="dash", line_color="red",
                  annotation_text=f"LCL: {lcl:.2f}%")
    
    fig.update_layout(
        title="p管理図（不良率管理図）",
        xaxis_title="日付",
        yaxis_title="不良率 (%)",
        height=500
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # 異常値の検出
    outliers = daily_summary[(daily_summary['不良率'] > ucl) | (daily_summary['不良率'] < lcl)]
    
    if len(outliers) > 0:
        st.warning(f"⚠️ {len(outliers)}件の異常値が検出されました")
        st.dataframe(outliers)
    else:
        st.success("✅ 異常値は検出されませんでした")

def create_checklist(df):
    """チェックシートの作成"""
    st.subheader("📋 チェックシート - データ収集支援")
    
    # 不良項目別のチェックリスト
    if '不良項目' in df.columns:
        st.write("**不良項目別チェックリスト**")
        
        defect_items = df['不良項目'].unique()
        
        for item in defect_items:
            with st.expander(f"📌 {item}"):
                item_data = df[df['不良項目'] == item]
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.write(f"**発生件数**: {len(item_data)}件")
                    st.write(f"**総不良数**: {item_data['不良数'].sum()}件")
                
                with col2:
                    if '原因分類' in df.columns:
                        causes = item_data['原因分類'].value_counts()
                        st.write("**主な原因**:")
                        for cause, count in causes.head(3).items():
                            st.write(f"• {cause}: {count}件")
    
    # カスタムチェックリスト
    st.write("**カスタムチェックリスト**")
    
    checklist_items = st.text_area(
        "チェック項目を入力してください（1行に1項目）",
        value="寸法測定\n表面粗さ検査\n外観検査\n機能テスト\n包装確認",
        height=100
    )
    
    if st.button("チェックリストを生成"):
        items = [item.strip() for item in checklist_items.split('\n') if item.strip()]
        
        st.write("**生成されたチェックシート**")
        for i, item in enumerate(items, 1):
            st.checkbox(f"{i}. {item}", key=f"check_{i}")

def create_graphs(df):
    """グラフの作成"""
    st.subheader("📊 グラフ - 時系列・比較分析")
    
    # グラフタイプの選択
    graph_type = st.selectbox(
        "グラフタイプを選択してください",
        ["時系列グラフ", "工程別比較", "製品別比較", "原因別比較"]
    )
    
    if graph_type == "時系列グラフ":
        if '日付' not in df.columns:
            st.error("データに「日付」の列が必要です。")
            return
        
        # 日付別の集計
        daily_summary = df.groupby('日付').agg({
            '不良数': 'sum',
            '検査数': 'sum'
        }).reset_index()
        daily_summary['不良率'] = (daily_summary['不良数'] / daily_summary['検査数'] * 100).round(2)
        
        # 時系列グラフ
        fig = px.line(daily_summary, x='日付', y='不良率',
                      title="不良率推移（時系列）")
        st.plotly_chart(fig, use_container_width=True)
    
    elif graph_type == "工程別比較":
        if '発生工程' not in df.columns:
            st.error("データに「発生工程」の列が必要です。")
            return
        
        # 工程別の集計
        process_summary = df.groupby('発生工程').agg({
            '不良数': 'sum',
            '検査数': 'sum'
        }).reset_index()
        process_summary['不良率'] = (process_summary['不良数'] / process_summary['検査数'] * 100).round(2)
        
        # 工程別比較グラフ
        fig = px.bar(process_summary, x='発生工程', y='不良率',
                     title="工程別不良率比較")
        st.plotly_chart(fig, use_container_width=True)
    
    elif graph_type == "製品別比較":
        if '製品名' not in df.columns:
            st.error("データに「製品名」の列が必要です。")
            return
        
        # 製品別の集計
        product_summary = df.groupby('製品名').agg({
            '不良数': 'sum',
            '検査数': 'sum'
        }).reset_index()
        product_summary['不良率'] = (product_summary['不良数'] / product_summary['検査数'] * 100).round(2)
        
        # 製品別比較グラフ
        fig = px.bar(product_summary, x='製品名', y='不良率',
                     title="製品別不良率比較")
        st.plotly_chart(fig, use_container_width=True)
    
    elif graph_type == "原因別比較":
        if '原因分類' not in df.columns:
            st.error("データに「原因分類」の列が必要です。")
            return
        
        # 原因別の集計
        cause_summary = df.groupby('原因分類')['不良数'].sum().sort_values(ascending=False)
        
        # 原因別比較グラフ
        fig = px.pie(values=cause_summary.values, names=cause_summary.index,
                     title="原因分類別不良数")
        st.plotly_chart(fig, use_container_width=True)


def generate_presentation(df):
    """プレゼン資料生成"""
    st.subheader("📊 プレゼン資料自動生成")
    
    # 会社情報入力
    st.write("**プレゼン資料の基本情報**")
    col1, col2 = st.columns(2)
    
    with col1:
        company_name = st.text_input("会社名", value="株式会社サンプル")
        department = st.text_input("部署名", value="品質管理部")
        presenter = st.text_input("発表者名", value="品質管理担当")
    
    with col2:
        presentation_date = st.date_input("発表日", value=datetime.now().date())
        period = st.text_input("分析期間", value="2024年1月")
        target = st.text_input("改善目標", value="不良率10%削減")
    
    # 生成するスライドの選択
    st.write("**生成するスライドを選択**")
    slide_options = {
        "📈 概要・データサマリー": True,
        "📊 パレート図分析": True,
        "🐟 特性要因図分析": True,
        "📈 管理図分析": True,
        "📊 時系列分析": True,
        "🎯 改善提案・アクションプラン": True
    }
    
    for option, default in slide_options.items():
        slide_options[option] = st.checkbox(option, value=default)
    
    # プレゼン資料生成ボタン
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("📊 PowerPoint生成"):
            pptx_file = create_pptx_presentation(df, company_name, department, presenter, 
                                               presentation_date, period, target, slide_options)
            if pptx_file:
                with open(pptx_file, "rb") as file:
                    st.download_button(
                        label="📊 PowerPointファイルをダウンロード",
                        data=file.read(),
                        file_name=f"QC7つ道具分析報告_{presentation_date.strftime('%Y%m%d')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
    
    with col2:
        if st.button("📄 PDF生成"):
            pdf_file = create_pdf_report(df, company_name, department, presenter, 
                                       presentation_date, period, target, slide_options)
            if pdf_file:
                with open(pdf_file, "rb") as file:
                    st.download_button(
                        label="📄 PDFファイルをダウンロード",
                        data=file.read(),
                        file_name=f"QC7つ道具分析報告_{presentation_date.strftime('%Y%m%d')}.pdf",
                        mime="application/pdf"
                    )

def create_pptx_presentation(df, company_name, department, presenter, 
                           presentation_date, period, target, slide_options):
    """PowerPointプレゼン資料作成"""
    try:
        # 新しいプレゼンテーション作成
        prs = Presentation()
        
        # スライド1: タイトルスライド
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = f"不良分析QC7つ道具\n分析報告書"
        subtitle.text = f"{company_name} {department}\n発表者: {presenter}\n{presentation_date.strftime('%Y年%m月%d日')}"
        
        # スライド2: 概要・データサマリー
        if slide_options["📈 概要・データサマリー"]:
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            summary_slide.shapes.title.text = "📈 分析概要・データサマリー"
            
            # データサマリーの計算
            total_defects = df['不良数'].sum()
            total_inspections = df['検査数'].sum()
            defect_rate = (total_defects / total_inspections * 100) if total_inspections > 0 else 0
            
            content = summary_slide.placeholders[1].text_frame
            content.text = f"""
分析期間: {period}
総検査数: {total_inspections:,}件
総不良数: {total_defects:,}件
不良率: {defect_rate:.2f}%
改善目標: {target}

主要な不良項目:
{', '.join(df.groupby('不良項目')['不良数'].sum().sort_values(ascending=False).head(3).index.tolist())}

分析目的:
• 不良原因の特定と対策立案
• 工程の安定性評価
• 改善活動の効果測定
"""
        
        # スライド3: パレート図分析
        if slide_options["📊 パレート図分析"]:
            pareto_slide = prs.slides.add_slide(prs.slide_layouts[1])
            pareto_slide.shapes.title.text = "📊 パレート図分析"
            
            # パレート図の作成
            defect_summary = df.groupby('不良項目')['不良数'].sum().sort_values(ascending=False)
            cumulative_ratio = (defect_summary.cumsum() / defect_summary.sum() * 100).round(1)
            
            content = pareto_slide.placeholders[1].text_frame
            content.text = f"""
重点管理項目（上位80%）:
{chr(10).join([f"• {item}: {count}件 ({ratio:.1f}%)" for item, count, ratio in 
              zip(defect_summary.index[:3], defect_summary.values[:3], 
                  [count/defect_summary.sum()*100 for count in defect_summary.values[:3]])])}

改善優先順位:
1. {defect_summary.index[0]}: {defect_summary.values[0]}件
2. {defect_summary.index[1]}: {defect_summary.values[1]}件  
3. {defect_summary.index[2]}: {defect_summary.values[2]}件

【改善提案】
上位3項目に集中した改善活動を実施し、
全体の不良率を10%以上削減する。
"""
        
        # スライド4: 特性要因図分析
        if slide_options["🐟 特性要因図分析"]:
            fishbone_slide = prs.slides.add_slide(prs.slide_layouts[1])
            fishbone_slide.shapes.title.text = "🐟 特性要因図分析"
            
            # 4M分析
            cause_summary = df.groupby('原因分類')['不良数'].sum().sort_values(ascending=False)
            
            content = fishbone_slide.placeholders[1].text_frame
            content.text = f"""
4M分析結果:

【Man（人）】作業者関連
• 取扱不注意: {df[df['原因分類'] == '作業者']['不良数'].sum()}件

【Machine（機械）】設備・工具関連  
• 加工不良: {df[df['原因分類'] == '加工']['不良数'].sum()}件
• 工具不良: {df[df['原因分類'] == '工具']['不良数'].sum()}件

【Material（材料）】材料関連
• 材料不良: {df[df['原因分類'] == '材料']['不良数'].sum()}件

【Method（方法）】環境・条件関連
• 環境要因: {df[df['原因分類'] == '環境']['不良数'].sum()}件

【改善方針】
原因別の体系的対策を立案し、
再発防止体制を構築する。
"""
        
        # スライド5: 管理図分析
        if slide_options["📈 管理図分析"]:
            control_slide = prs.slides.add_slide(prs.slide_layouts[1])
            control_slide.shapes.title.text = "📈 管理図分析"
            
            # 日付別の集計
            daily_summary = df.groupby('日付').agg({
                '不良数': 'sum',
                '検査数': 'sum'
            }).reset_index()
            daily_summary['不良率'] = (daily_summary['不良数'] / daily_summary['検査数'] * 100).round(2)
            
            p_bar = daily_summary['不良率'].mean()
            n_bar = daily_summary['検査数'].mean()
            ucl = p_bar + 3 * np.sqrt(p_bar * (100 - p_bar) / n_bar)
            lcl = max(0, p_bar - 3 * np.sqrt(p_bar * (100 - p_bar) / n_bar))
            
            outliers = daily_summary[(daily_summary['不良率'] > ucl) | (daily_summary['不良率'] < lcl)]
            
            content = control_slide.placeholders[1].text_frame
            content.text = f"""
工程管理状況:

中心線: {p_bar:.2f}%
UCL: {ucl:.2f}%
LCL: {lcl:.2f}%

異常値: {len(outliers)}件検出
{'・' + chr(10) + '・'.join(outliers['日付'].astype(str).tolist()) if len(outliers) > 0 else '異常値なし'}

【工程評価】
{'工程は不安定。異常値の原因調査が必要。' if len(outliers) > 0 else '工程は比較的安定。継続的な監視が必要。'}

【改善方針】
管理限界を超えた日の原因を詳細調査し、
工程の安定化を図る。
"""
        
        # スライド6: 時系列分析
        if slide_options["📊 時系列分析"]:
            timeseries_slide = prs.slides.add_slide(prs.slide_layouts[1])
            timeseries_slide.shapes.title.text = "📊 時系列分析"
            
            # 日付別トレンド
            daily_trend = daily_summary['不良率'].values
            trend_direction = "改善" if daily_trend[-1] < daily_trend[0] else "悪化" if daily_trend[-1] > daily_trend[0] else "横ばい"
            
            content = timeseries_slide.placeholders[1].text_frame
            content.text = f"""
不良率推移:

期間開始時: {daily_trend[0]:.2f}%
期間終了時: {daily_trend[-1]:.2f}%
変化: {trend_direction} ({daily_trend[-1] - daily_trend[0]:+.2f}%)

最高値: {daily_trend.max():.2f}%
最低値: {daily_trend.min():.2f}%
平均値: {daily_trend.mean():.2f}%

【トレンド分析】
{'改善傾向が見られる。' if trend_direction == "改善" else '悪化傾向。緊急対策が必要。' if trend_direction == "悪化" else '横ばい傾向。改善活動の効果が限定的。'}

【今後の方針】
継続的な改善活動により、
目標の{target}を達成する。
"""
        
        # スライド7: 改善提案・アクションプラン
        if slide_options["🎯 改善提案・アクションプラン"]:
            action_slide = prs.slides.add_slide(prs.slide_layouts[1])
            action_slide.shapes.title.text = "🎯 改善提案・アクションプラン"
            
            content = action_slide.placeholders[1].text_frame
            content.text = f"""
改善提案:

【短期対策（1ヶ月以内）】
• 重点不良項目の原因調査強化
• 作業標準の見直しと徹底
• 検査頻度の増加

【中期対策（3ヶ月以内）】
• 設備・工具の点検体制強化
• 作業者教育の充実
• 工程能力の向上

【長期対策（6ヶ月以内）】
• 工程設計の見直し
• 自動化・省人化の検討
• 品質管理システムの構築

【期待効果】
不良率: {defect_rate:.2f}% → {defect_rate * 0.9:.2f}% (10%削減)
コスト削減: 年間約XXX万円
品質向上: 顧客満足度向上

【実施体制】
責任者: {presenter}
期間: {period} - {period}
次回レビュー: {(presentation_date + timedelta(days=30)).strftime('%Y年%m月%d日')}
"""
        
        # ファイル保存
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        return temp_file.name
        
    except Exception as e:
        st.error(f"PowerPoint生成エラー: {str(e)}")
        return None

def create_pdf_report(df, company_name, department, presenter, 
                     presentation_date, period, target, slide_options):
    """PDFレポート作成"""
    try:
        # PDFファイル作成
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        doc = SimpleDocTemplate(temp_file.name, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # タイトル
        title = Paragraph(f"不良分析QC7つ道具 分析報告書", styles['Title'])
        story.append(title)
        story.append(Spacer(1, 12))
        
        # 基本情報
        info_text = f"""
<para align=center>
<b>{company_name} {department}</b><br/>
発表者: {presenter}<br/>
発表日: {presentation_date.strftime('%Y年%m月%d日')}<br/>
分析期間: {period}
</para>
"""
        story.append(Paragraph(info_text, styles['Normal']))
        story.append(Spacer(1, 20))
        
        # データサマリー
        if slide_options["📈 概要・データサマリー"]:
            total_defects = df['不良数'].sum()
            total_inspections = df['検査数'].sum()
            defect_rate = (total_defects / total_inspections * 100) if total_inspections > 0 else 0
            
            summary_text = f"""
<h2>📈 分析概要・データサマリー</h2>
<p>
総検査数: {total_inspections:,}件<br/>
総不良数: {total_defects:,}件<br/>
不良率: {defect_rate:.2f}%<br/>
改善目標: {target}<br/>
</p>
"""
            story.append(Paragraph(summary_text, styles['Normal']))
            story.append(Spacer(1, 12))
        
        # パレート図分析
        if slide_options["📊 パレート図分析"]:
            defect_summary = df.groupby('不良項目')['不良数'].sum().sort_values(ascending=False)
            
            pareto_text = f"""
<h2>📊 パレート図分析</h2>
<p>
<b>重点管理項目（上位3位）:</b><br/>
1. {defect_summary.index[0]}: {defect_summary.values[0]}件<br/>
2. {defect_summary.index[1]}: {defect_summary.values[1]}件<br/>
3. {defect_summary.index[2]}: {defect_summary.values[2]}件<br/>
</p>
"""
            story.append(Paragraph(pareto_text, styles['Normal']))
            story.append(Spacer(1, 12))
        
        # 特性要因図分析
        if slide_options["🐟 特性要因図分析"]:
            cause_summary = df.groupby('原因分類')['不良数'].sum().sort_values(ascending=False)
            
            fishbone_text = f"""
<h2>🐟 特性要因図分析</h2>
<p>
<b>4M分析結果:</b><br/>
• Man（人）: {df[df['原因分類'] == '作業者']['不良数'].sum()}件<br/>
• Machine（機械）: {df[df['原因分類'].isin(['加工', '工具'])]['不良数'].sum()}件<br/>
• Material（材料）: {df[df['原因分類'] == '材料']['不良数'].sum()}件<br/>
• Method（方法）: {df[df['原因分類'] == '環境']['不良数'].sum()}件<br/>
</p>
"""
            story.append(Paragraph(fishbone_text, styles['Normal']))
            story.append(Spacer(1, 12))
        
        # 改善提案
        if slide_options["🎯 改善提案・アクションプラン"]:
            action_text = f"""
<h2>🎯 改善提案・アクションプラン</h2>
<p>
<b>短期対策（1ヶ月以内）:</b><br/>
• 重点不良項目の原因調査強化<br/>
• 作業標準の見直しと徹底<br/>
• 検査頻度の増加<br/>
</p>
<p>
<b>中期対策（3ヶ月以内）:</b><br/>
• 設備・工具の点検体制強化<br/>
• 作業者教育の充実<br/>
• 工程能力の向上<br/>
</p>
<p>
<b>期待効果:</b><br/>
不良率10%削減、年間コスト削減XXX万円、品質向上<br/>
</p>
"""
            story.append(Paragraph(action_text, styles['Normal']))
        
        # PDF生成
        doc.build(story)
        return temp_file.name
        
    except Exception as e:
        st.error(f"PDF生成エラー: {str(e)}")
        return None

if __name__ == "__main__":
    pass
